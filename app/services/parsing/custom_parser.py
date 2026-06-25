import asyncio
import csv
import io
from typing import Any, BinaryIO, Dict, List, Optional

from app.services.parsing.base import BaseDocumentParser, PageContent, ParsedDocument

_SUPPORTED = {
    "application/pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "text/csv",
    "text/plain",
}


class CustomParser(BaseDocumentParser):
    def __init__(self, ocr_engine: str = "rapidocr", ocr_languages: List[str] = None):
        self.ocr_engine = ocr_engine
        self.ocr_languages = ocr_languages or ["en"]

    def supports(self, mime_type: str) -> bool:
        return mime_type in _SUPPORTED

    async def parse(
        self,
        file: BinaryIO,
        mime_type: str,
        filename: str,
        options: Optional[Dict[str, Any]] = None,
    ) -> ParsedDocument:
        # options apply to docling-serve only; the custom parser ignores them.
        return await asyncio.to_thread(self._parse_sync, file, mime_type, filename)

    def _parse_sync(self, file: BinaryIO, mime_type: str, filename: str) -> ParsedDocument:
        if mime_type == "application/pdf":
            return self._parse_pdf(file)
        if mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            return self._parse_docx(file)
        if mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            return self._parse_pptx(file)
        if mime_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
            return self._parse_xlsx(file)
        if mime_type in ("text/csv", "text/plain"):
            return self._parse_text(file, mime_type)
        raise ValueError(f"Unsupported mime type: {mime_type!r}")

    def _parse_pdf(self, file: BinaryIO) -> ParsedDocument:
        from pypdf import PdfReader
        reader = PdfReader(file)
        pages = []
        all_text = []
        for page_no, page in enumerate(reader.pages, start=1):
            text = page.extract_text() or ""
            all_text.append(text)
            pages.append(PageContent(page_no=page_no, text=text))
        return ParsedDocument(
            text="\n\n".join(all_text),
            pages=pages,
            metadata={"page_count": len(pages)},
        )

    def _parse_docx(self, file: BinaryIO) -> ParsedDocument:
        import docx
        doc = docx.Document(file)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return ParsedDocument(text="\n\n".join(paragraphs), metadata={})

    def _parse_pptx(self, file: BinaryIO) -> ParsedDocument:
        from pptx import Presentation
        prs = Presentation(file)
        pages = []
        all_text = []
        for slide_no, slide in enumerate(prs.slides, start=1):
            slide_text = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            slide_text.append(para.text)
            text = "\n".join(slide_text)
            all_text.append(text)
            pages.append(PageContent(page_no=slide_no, text=text))
        return ParsedDocument(
            text="\n\n".join(all_text),
            pages=pages,
            metadata={"slide_count": len(pages)},
        )

    def _parse_xlsx(self, file: BinaryIO) -> ParsedDocument:
        import openpyxl
        wb = openpyxl.load_workbook(file, data_only=True)
        all_text = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                row_text = "\t".join(str(v) if v is not None else "" for v in row)
                if row_text.strip():
                    rows.append(row_text)
            if rows:
                all_text.append(f"## Sheet: {sheet_name}\n" + "\n".join(rows))
        return ParsedDocument(
            text="\n\n".join(all_text),
            metadata={"sheet_count": len(wb.sheetnames)},
        )

    def _parse_text(self, file: BinaryIO, mime_type: str) -> ParsedDocument:
        raw = file.read()
        text = raw.decode("utf-8", errors="replace")
        if mime_type == "text/csv":
            reader = csv.reader(io.StringIO(text))
            text = "\n".join("\t".join(row) for row in reader)
        return ParsedDocument(text=text, metadata={})
